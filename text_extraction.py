import pandas as pd
import docx
from PyPDF2 import PdfReader
from pptx import Presentation
from txtai.pipeline import Labels

def read_csv(file_path):
    try:
        df = pd.read_csv(file_path)
        texts = df.to_string()  # Convert the entire DataFrame to a string
        print("CSV content:")
        print(texts)
        return texts
    except Exception as e:
        print(f"Error reading CSV: {e}")

def read_docx(file_path):
    try:
        doc = docx.Document(file_path)
        text = [paragraph.text for paragraph in doc.paragraphs]
        full_text = "\n".join(text)
        print("DOCX content:")
        print(full_text)
        return full_text
    except Exception as e:
        print(f"Error reading DOCX: {e}")

def read_pdf(file_path):
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            text = []
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text.append(page.extract_text())
        full_text = "\n".join(text)
        print("PDF content:")
        print(full_text)
        return full_text
    except Exception as e:
        print(f"Error reading PDF: {e}")

def read_ppt(file_path):
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        full_text = "\n".join(text)
        print("PPT content:")
        print(full_text)
        return full_text
    except Exception as e:
        print(f"Error reading PPT: {e}")

def read_txt(file_path):
    try:
        with open(file_path, 'r') as file:
            text = file.read()
        print("TXT content:")
        print(text)
        return text
    except Exception as e:
        print(f"Error reading TXT: {e}")

# Test CSV
csv_text = read_csv('Docs/FAQ_Chatbot.csv')

# Test DOCX
docx_text = read_docx('Docs/LADP #2 Question Bank_.docx')

# Test PDF
pdf_text = read_pdf('Docs/Guide to Right Classification - NYP ppt v190801.pdf')

# Test PPT
ppt_text = read_ppt('Docs/Intake 2 LADP Knowledge Transfer 1 Slides.pptx')

# Test TXT
txt_text = read_txt('Docs/sgpresident.txt')

# Example usage of txtai (if needed for later processing)
# Create a labels pipeline
labels = Labels()

# Example text for classification
example_text = "This is a confidential document."

# Example labels
labelset = ["Top Secret", "Confidential", "Restricted", "Official"]

# Apply labels
prediction = labels(example_text, labelset)
print(f"Text: {example_text}\nPrediction: {prediction}\n")




#txtai test
#from txtai.pipeline import Textractor
#
# # Create textractor model
# textractor = Textractor()
# text = textractor("Docs/Guide to Right Classification - NYP ppt v190801.pdf")
# print(text)
